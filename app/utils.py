import os
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import SentenceTransformerEmbeddings
from langchain.vectorstores import FAISS
from langchain.llms import Ollama
from langchain.chains import ConversationalRetrievalChain
from langchain.memory import ConversationBufferMemory
from langchain.schema import Document

# Maximum number of tokens supported by the model.
MAX_TOKENS = 8192

# Default model used by the Ollama API.
DEFAULT_MODEL = os.getenv("OLLAMA_MODEL", "llama3")


def summarize_text(text: str, model: str = DEFAULT_MODEL) -> str:
    """Return a comprehensive summary of ``text`` using ``model``.

    The text is processed in 4000 character chunks to ensure the model has
    sufficient context. Each chunk is summarized individually and the
    resulting summaries are combined into a final overview.
    """

    if not text:
        return ""

    base_url = os.getenv("OLLAMA_BASE_URL", "http://ollama:11434")
    llm = Ollama(model=model, base_url=base_url, temperature=0.1)

    chunk_size = 4000
    chunks = [text[i : i + chunk_size] for i in range(0, len(text), chunk_size)]
    partial_summaries: list[str] = []

    for chunk in chunks:
        prompt = f"Summarize the following text:\n\n{chunk}"
        partial_summaries.append(llm.invoke(prompt).strip())

    if len(partial_summaries) == 1:
        return partial_summaries[0]

    combo_prompt = (
        "Combine the following summaries into a single comprehensive summary:\n\n"
        + "\n".join(partial_summaries)
    )
    return llm.invoke(combo_prompt).strip()


def load_docx(path: str) -> list[Document]:
    """Return the text content from a DOCX file."""
    doc = docx.Document(path)
    docs = []
    for i, paragraph in enumerate(doc.paragraphs, start=1):
        if paragraph.text.strip():
            docs.append(Document(page_content=paragraph.text, metadata={"paragraph": i}))
    return docs


def load_txt(path: str) -> list[Document]:
    """Return the text content from a plain text file.

    The function first attempts to decode the file as UTF-8. If that fails
    due to a ``UnicodeDecodeError`` it falls back to a list of common
    encodings before finally ignoring undecodable bytes. This prevents
    crashes when the text file uses a legacy code page such as Windows-1252.
    """

    with open(path, "rb") as f:
        raw = f.read()

    encodings = ["utf-8", "cp1252", "latin-1"]
    for enc in encodings:
        try:
            text = raw.decode(enc)
            return [Document(page_content=text)]
        except UnicodeDecodeError:
            continue

    # As a last resort, drop undecodable bytes.
    text = raw.decode("utf-8", errors="ignore")
    return [Document(page_content=text)]


def load_pptx(path: str) -> list[Document]:
    """Return the text content from a PowerPoint file."""
    prs = Presentation(path)
    docs = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)
        if slide_text:
            docs.append(
                Document(page_content="\n".join(slide_text), metadata={"page": i})
            )
    return docs


def load_pdf_with_password(path: str, password: str) -> list[Document]:
    """Return the text from a PDF, unlocking it with the given password."""
    with fitz.open(path) as doc:
        if doc.needs_pass and not doc.authenticate(password):
            raise ValueError("Incorrect password")
        docs = [
            Document(page_content=page.get_text(), metadata={"page": i + 1})
            for i, page in enumerate(doc)
        ]
    return docs


def load_document(path: str, password: str = "") -> list[Document]:
    """Load text from various document formats based on file extension."""
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return load_pdf_with_password(path, password)
    if ext == ".docx":
        return load_docx(path)
    if ext == ".txt":
        return load_txt(path)
    if ext in {".ppt", ".pptx"}:
        return load_pptx(path)
    return []


def create_vector_store(docs: list[Document], path: str) -> FAISS:
    """Create (or load) a persistent FAISS vector store from ``docs``."""
    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    chunks = splitter.split_documents(docs)
    embeddings = SentenceTransformerEmbeddings(model_name="all-MiniLM-L6-v2")
    if os.path.exists(path):
        return FAISS.load_local(
            path,
            embeddings,
            allow_dangerous_deserialization=True,
        )
    vector_store = FAISS.from_documents(chunks, embeddings)
    vector_store.save_local(path)
    return vector_store


def get_qa_chain(
    vector_store: FAISS,
    temperature: float = 0.1,
    max_tokens: int = 8192,
    model: str = DEFAULT_MODEL,
) -> ConversationalRetrievalChain:
    """Return a conversational QA chain with memory using ``model``."""
    base_url = os.getenv("OLLAMA_BASE_URL", "http://ollama:11434")
    llm = Ollama(
        model=model,
        base_url=base_url,
        temperature=temperature,
        num_predict=min(max_tokens, MAX_TOKENS),
    )
    memory = ConversationBufferMemory(
        memory_key="chat_history",
        return_messages=True,
        output_key="answer",
    )
    return ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vector_store.as_retriever(),
        memory=memory,
        return_source_documents=True,
    )


def rewrite_question(
    question: str, history: list[tuple[str, str]] | None = None,
    model: str = DEFAULT_MODEL,
) -> str:
    """Return a clarified version of ``question`` using conversation context."""
    base_url = os.getenv("OLLAMA_BASE_URL", "http://ollama:11434")
    llm = Ollama(model=model, base_url=base_url, temperature=0.0)
    history_text = ""
    if history:
        history_text = "\n".join(f"{r}: {m}" for r, m in history[-6:])
    prompt = (
        "Rewrite the following user question to be clear and self contained."\
        " Use the prior conversation for context if helpful.\n\n"\
        f"History:\n{history_text}\n\nQuestion: {question}\n\nRewritten:"
    )
    try:
        return llm.invoke(prompt).strip()
    except Exception:
        return question

